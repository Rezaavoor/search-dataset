"""Format-specific page extractors for the document page store.

Each loader returns a list of (page_number, text_content, metadata_dict) tuples.
Page numbers are 1-indexed.  For formats without natural page boundaries
(DOCX, TXT, DOC), text is split into ~3000-character chunks that approximate a
typical PDF page of legal text.

Supported formats:
    pdf, docx, xlsx, pptx, txt, csv, json, xls, doc
"""

import csv
import io
import json
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# ---------------------------------------------------------------------------
# Type alias for a single extracted "page"
# ---------------------------------------------------------------------------
PageTuple = Tuple[int, str, Dict[str, Any]]
# (page_number, text_content, extra_metadata)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
TEXT_CHUNK_SIZE = 3000          # chars – approximate one PDF-page equivalent
TEXT_CHUNK_OVERLAP = 200        # chars – overlap between chunks
MAX_SHEET_CHARS = 5000          # split large spreadsheet sheets at this limit
MAX_CSV_ROWS_PER_PAGE = 50     # for CSV chunking

# Supported extensions (lowercase, with leading dot)
SUPPORTED_EXTENSIONS: Dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".xls": "xls",
    ".pptx": "pptx",
    ".txt": "txt",
    ".csv": "csv",
    ".json": "json",
    ".doc": "doc",
}


def file_type_from_path(path: Path) -> Optional[str]:
    """Return the canonical file_type string for a path, or None if unsupported."""
    return SUPPORTED_EXTENSIONS.get(path.suffix.lower())


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
def load_file_pages(path: Path) -> Tuple[str, List[PageTuple]]:
    """Route to format-specific loader.

    Returns (file_type, list_of_page_tuples).
    Raises ValueError for unsupported extensions.
    """
    ext = path.suffix.lower()
    ft = SUPPORTED_EXTENSIONS.get(ext)
    if ft is None:
        raise ValueError(f"Unsupported file extension: {ext!r} ({path.name})")

    loader = _LOADERS.get(ft)
    if loader is None:
        raise ValueError(f"No loader registered for file_type={ft!r}")
    pages = loader(path)
    # Sanity: ensure page_numbers are 1-indexed and sequential
    for i, (pn, _txt, _md) in enumerate(pages):
        if pn != i + 1:
            # Re-number if the loader produced inconsistent numbering
            pages = [(i + 1, txt, md) for i, (_, txt, md) in enumerate(pages)]
            break
    return ft, pages


# ---------------------------------------------------------------------------
# Text chunking helper (used by DOCX, TXT, DOC, and oversized sheets)
# ---------------------------------------------------------------------------
def _chunk_text(
    text: str,
    *,
    chunk_size: int = TEXT_CHUNK_SIZE,
    chunk_overlap: int = TEXT_CHUNK_OVERLAP,
) -> List[str]:
    """Split *text* into roughly equal chunks on whitespace boundaries.

    If the text is shorter than *chunk_size* it's returned as a single chunk.
    """
    text = text.strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]

    chunks: List[str] = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end >= len(text):
            chunks.append(text[start:].strip())
            break
        # Try to break at whitespace near the chunk boundary
        candidate = text[start:end]
        last_space = candidate.rfind(" ")
        if last_space > chunk_size // 2:
            end = start + last_space
        chunks.append(text[start:end].strip())
        start = end - chunk_overlap
        if start < 0:
            start = 0

    # Remove empty trailing chunks
    return [c for c in chunks if c]


# ---------------------------------------------------------------------------
# PDF loader
# ---------------------------------------------------------------------------
def _load_pdf(path: Path) -> List[PageTuple]:
    from langchain_community.document_loaders import PyPDFLoader

    loader = PyPDFLoader(str(path))
    docs = loader.load()
    docs = sorted(
        docs,
        key=lambda d: d.metadata.get("page")
        if isinstance(d.metadata.get("page"), int)
        else 10**9,
    )
    pages: List[PageTuple] = []
    for idx, doc in enumerate(docs):
        page0 = doc.metadata.get("page")
        page_number = int(page0) + 1 if isinstance(page0, int) else (idx + 1)
        content = str(doc.page_content or "")
        md = dict(doc.metadata or {}) if isinstance(doc.metadata, dict) else {}
        md["page"] = page_number - 1  # keep 0-indexed in metadata (LangChain convention)
        pages.append((page_number, content, md))
    return pages


# ---------------------------------------------------------------------------
# DOCX loader
# ---------------------------------------------------------------------------
def _load_docx(path: Path) -> List[PageTuple]:
    import docx2txt

    text = docx2txt.process(str(path))
    if not text or not text.strip():
        return [(1, "", {"format_note": "empty_docx"})]

    chunks = _chunk_text(text)
    pages: List[PageTuple] = []
    for idx, chunk in enumerate(chunks):
        pages.append((
            idx + 1,
            chunk,
            {"page": idx, "chunk_method": "text_split_3k"},
        ))
    return pages


# ---------------------------------------------------------------------------
# XLSX loader
# ---------------------------------------------------------------------------
def _load_xlsx(path: Path) -> List[PageTuple]:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages: List[PageTuple] = []
    page_num = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text: List[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            # Skip completely empty rows
            if not any(c.strip() for c in cells):
                continue
            rows_text.append(" | ".join(cells))

        if not rows_text:
            continue

        sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text)

        # If the sheet text is too long, chunk it
        if len(sheet_text) > MAX_SHEET_CHARS:
            sub_chunks = _chunk_text(sheet_text, chunk_size=MAX_SHEET_CHARS)
            for sub_idx, chunk in enumerate(sub_chunks):
                page_num += 1
                pages.append((
                    page_num,
                    chunk,
                    {
                        "page": page_num - 1,
                        "sheet_name": sheet_name,
                        "sheet_part": sub_idx + 1,
                        "chunk_method": "sheet_split",
                    },
                ))
        else:
            page_num += 1
            pages.append((
                page_num,
                sheet_text,
                {"page": page_num - 1, "sheet_name": sheet_name},
            ))

    wb.close()

    if not pages:
        return [(1, "", {"format_note": "empty_xlsx"})]
    return pages


# ---------------------------------------------------------------------------
# XLS loader (old Excel format)
# ---------------------------------------------------------------------------
def _load_xls(path: Path) -> List[PageTuple]:
    """Load old-format .xls via pandas (requires xlrd)."""
    try:
        import pandas as pd

        xls = pd.ExcelFile(str(path), engine="xlrd")
    except ImportError:
        # Fall back: try openpyxl (won't work for real .xls but worth trying)
        return _load_xlsx(path)
    except Exception:
        return [(1, f"[Could not read XLS file: {path.name}]", {"format_note": "xls_read_error"})]

    pages: List[PageTuple] = []
    page_num = 0
    for sheet_name in xls.sheet_names:
        try:
            df = pd.read_excel(xls, sheet_name=sheet_name)
        except Exception:
            continue
        if df.empty:
            continue
        sheet_text = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
        if len(sheet_text) > MAX_SHEET_CHARS:
            sub_chunks = _chunk_text(sheet_text, chunk_size=MAX_SHEET_CHARS)
            for sub_idx, chunk in enumerate(sub_chunks):
                page_num += 1
                pages.append((
                    page_num,
                    chunk,
                    {
                        "page": page_num - 1,
                        "sheet_name": sheet_name,
                        "sheet_part": sub_idx + 1,
                        "chunk_method": "sheet_split",
                    },
                ))
        else:
            page_num += 1
            pages.append((
                page_num,
                sheet_text,
                {"page": page_num - 1, "sheet_name": sheet_name},
            ))
    xls.close()

    if not pages:
        return [(1, "", {"format_note": "empty_xls"})]
    return pages


# ---------------------------------------------------------------------------
# PPTX loader
# ---------------------------------------------------------------------------
def _load_pptx(path: Path) -> List[PageTuple]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: List[PageTuple] = []

    for slide_idx, slide in enumerate(prs.slides):
        texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            # Also capture table content
            if shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))

        slide_text = "\n".join(texts)
        page_number = slide_idx + 1
        pages.append((
            page_number,
            slide_text,
            {"page": slide_idx, "slide_number": page_number},
        ))

    if not pages:
        return [(1, "", {"format_note": "empty_pptx"})]
    return pages


# ---------------------------------------------------------------------------
# TXT loader
# ---------------------------------------------------------------------------
def _load_txt(path: Path) -> List[PageTuple]:
    # Try common encodings
    text = None
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            text = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    if text is None:
        # Last resort: read as bytes, replace errors
        text = path.read_bytes().decode("utf-8", errors="replace")

    if not text.strip():
        return [(1, "", {"format_note": "empty_txt"})]

    chunks = _chunk_text(text)
    pages: List[PageTuple] = []
    for idx, chunk in enumerate(chunks):
        pages.append((
            idx + 1,
            chunk,
            {"page": idx, "chunk_method": "text_split_3k"},
        ))
    return pages


# ---------------------------------------------------------------------------
# CSV loader
# ---------------------------------------------------------------------------
def _load_csv(path: Path) -> List[PageTuple]:
    # Try common encodings
    raw = None
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            raw = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    if raw is None:
        raw = path.read_bytes().decode("utf-8", errors="replace")

    if not raw.strip():
        return [(1, "", {"format_note": "empty_csv"})]

    # Detect dialect
    try:
        dialect = csv.Sniffer().sniff(raw[:4096])
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(raw), dialect)
    all_rows = list(reader)
    if not all_rows:
        return [(1, "", {"format_note": "empty_csv"})]

    # First row as header
    header = all_rows[0] if all_rows else []
    data_rows = all_rows[1:] if len(all_rows) > 1 else all_rows

    pages: List[PageTuple] = []
    page_num = 0

    for chunk_start in range(0, max(len(data_rows), 1), MAX_CSV_ROWS_PER_PAGE):
        chunk_rows = data_rows[chunk_start : chunk_start + MAX_CSV_ROWS_PER_PAGE]
        lines = []
        if header:
            lines.append(" | ".join(header))
            lines.append("-" * min(len(" | ".join(header)), 80))
        for row in chunk_rows:
            lines.append(" | ".join(row))
        text = "\n".join(lines)
        page_num += 1
        pages.append((
            page_num,
            text,
            {
                "page": page_num - 1,
                "row_range": f"{chunk_start + 1}-{chunk_start + len(chunk_rows)}",
                "chunk_method": "csv_rows",
            },
        ))

    return pages if pages else [(1, "", {"format_note": "empty_csv"})]


# ---------------------------------------------------------------------------
# JSON loader
# ---------------------------------------------------------------------------
def _load_json(path: Path) -> List[PageTuple]:
    # Try common encodings
    raw = None
    for encoding in ("utf-8", "latin-1"):
        try:
            raw = path.read_text(encoding=encoding)
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    if raw is None:
        raw = path.read_bytes().decode("utf-8", errors="replace")

    if not raw.strip():
        return [(1, "", {"format_note": "empty_json"})]

    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        # Store the raw text as-is if it can't be parsed
        return [(1, raw[:TEXT_CHUNK_SIZE * 2], {"format_note": "invalid_json"})]

    pretty = json.dumps(data, indent=2, ensure_ascii=False, default=str)

    if len(pretty) <= TEXT_CHUNK_SIZE * 2:
        return [(1, pretty, {"page": 0})]

    # For large JSON, chunk the pretty-printed output
    chunks = _chunk_text(pretty, chunk_size=TEXT_CHUNK_SIZE)
    pages: List[PageTuple] = []
    for idx, chunk in enumerate(chunks):
        pages.append((
            idx + 1,
            chunk,
            {"page": idx, "chunk_method": "json_text_split"},
        ))
    return pages


# ---------------------------------------------------------------------------
# DOC loader (old Word format)
# ---------------------------------------------------------------------------
def _load_doc(path: Path) -> List[PageTuple]:
    """Best-effort loader for legacy .doc files.

    Strategy:
    1. Try textract if installed
    2. Try antiword if installed
    3. Fall back to reading raw bytes and extracting printable text
    """
    text = None

    # Strategy 1: textract
    try:
        import textract
        text = textract.process(str(path)).decode("utf-8", errors="replace")
    except Exception:
        pass

    # Strategy 2: antiword via subprocess
    if text is None:
        try:
            import subprocess
            result = subprocess.run(
                ["antiword", str(path)],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                text = result.stdout
        except Exception:
            pass

    # Strategy 3: raw byte extraction (very lossy)
    if text is None:
        try:
            raw = path.read_bytes()
            # Extract printable ASCII sequences of length >= 4
            printable = re.findall(rb"[\x20-\x7e]{4,}", raw)
            text = "\n".join(b.decode("ascii") for b in printable)
        except Exception:
            text = ""

    if not text or not text.strip():
        return [(1, "", {"format_note": "empty_or_unreadable_doc"})]

    chunks = _chunk_text(text)
    pages: List[PageTuple] = []
    for idx, chunk in enumerate(chunks):
        pages.append((
            idx + 1,
            chunk,
            {"page": idx, "chunk_method": "text_split_3k"},
        ))
    return pages


# ---------------------------------------------------------------------------
# Loader registry
# ---------------------------------------------------------------------------
_LOADERS = {
    "pdf": _load_pdf,
    "docx": _load_docx,
    "xlsx": _load_xlsx,
    "xls": _load_xls,
    "pptx": _load_pptx,
    "txt": _load_txt,
    "csv": _load_csv,
    "json": _load_json,
    "doc": _load_doc,
}
